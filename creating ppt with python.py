from pptx import Presentation

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide Layout
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Using For Loops with Dictionaries in Python"
subtitle.text = "Operations on Dictionaries"

# Slide 2: For Loops with Dictionaries
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Content Slide Layout
title = slide2.shapes.title
content = slide2.placeholders[1]
title.text = "For Loops with Dictionaries"

# Content for slide 2
content_text = """
- In Python, you can use for loops to iterate over dictionaries.
- You can iterate over keys, values, or key-value pairs.
- Example:
    for key in my_dict:
        print(key)
"""
content.text = content_text

# Slide 3: Operations on Dictionaries
slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Content Slide Layout
title = slide3.shapes.title
content = slide3.placeholders[1]
title.text = "Operations on Dictionaries"

# Content for slide 3
content_text = """
- Dictionaries are versatile data structures in Python.
- You can perform various operations on dictionaries:
    - Adding and updating key-value pairs.
    - Removing key-value pairs.
    - Checking if a key exists.
    - Iterating over keys, values, or items.
"""
content.text = content_text

# Save the PowerPoint presentation
prs.save("python_dict_operations.pptx")

print("Presentation created successfully.")
